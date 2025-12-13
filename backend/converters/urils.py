# utils.py
# Conversion utilities for PDF <-> other formats (pptx, docx, xlsx, images, etc.)

# === Standard Library Imports ===
import io
import os
import platform
import shutil
import subprocess
import tempfile
import zipfile
import logging
from io import BytesIO
from typing import List, Optional, Tuple

# === Third-Party Libraries ===
import pandas as pd
import pdfplumber
from pdf2image import convert_from_bytes
from PIL import Image, UnidentifiedImageError
import pytesseract

# PDFMiner (text extraction)
from pdfminer.high_level import extract_text as pdfminer_extract_text
from pdfminer.layout import LAParams

# python-pptx
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Length

from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas

# Optional: ReportLab (Excel fallback, TXT->PDF rendering)
try:
    from reportlab.lib.pagesizes import A4, letter
    from reportlab.platypus import (
        SimpleDocTemplate,
        Paragraph,
        Spacer,
        Table,
        TableStyle,
        PageBreak,
    )
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.pdfgen import canvas  # needed for TXT -> PDF

    HAS_REPORTLAB = True
except Exception:
    HAS_REPORTLAB = False

# Optional: pdf2docx (PDF -> DOCX)
try:
    from pdf2docx import Converter
    HAS_PDF2DOCX = True
except Exception:
    HAS_PDF2DOCX = False


logger = logging.getLogger(__name__)

# Configuration
_KEEP_TMP_ON_ERROR = True

# -------------------------
# TXT -> PDF
# -------------------------

# -------------------------
# TXT -> PDF (Improved, fixes black squares)
# -------------------------

def normalize_unicode_spaces(text: str) -> str:
    # Replace ONLY unicode special spaces, keep newlines and tabs as is
    special_spaces = [
        "\u00A0",  # NO-BREAK SPACE
        "\u2000", "\u2001", "\u2002", "\u2003", "\u2004", "\u2005",
        "\u2006", "\u2007", "\u2008", "\u2009", "\u200A",
        "\u202F", "\u205F", "\u3000",
    ]
    for sp in special_spaces:
        text = text.replace(sp, " ")
    return text


def txt_to_pdf_bytes(file_obj, page_size=A4, font_size=12, margin=40):
    try:
        file_obj.seek(0)
    except:
        pass

    raw = file_obj.read()
    try:
        text = raw.decode("utf-8")
    except:
        text = raw.decode("latin-1")

    # Normalize only special unicode spaces
    text = normalize_unicode_spaces(text)

    # Convert tabs → spaces (fixes black squares)
    text = text.replace("\t", "    ")  # 4 spaces (change to 8 if you want)

    # Register Unicode monospace font
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    pdfmetrics.registerFont(
        TTFont("DejaVuMono", "/usr/share/fonts/truetype/dejavu/DejaVuSansMono.ttf")
    )

    buffer = BytesIO()
    pdf = canvas.Canvas(buffer, pagesize=page_size)
    pdf.setFont("DejaVuMono", font_size)

    width, height = page_size
    x = margin
    y = height - margin
    line_height = font_size * 1.35

    for line in text.split("\n"):
        if y < margin:
            pdf.showPage()
            pdf.setFont("DejaVuMono", font_size)
            y = height - margin

        pdf.drawString(x, y, line)
        y -= line_height

    pdf.save()
    buffer.seek(0)
    return buffer.read()


# -------------------------
# Helper subprocess wrappers
# -------------------------

def _is_soffice_available() -> bool:
    return shutil.which("soffice") is not None


def _is_libreoffice_available() -> bool:
    return _is_soffice_available()


def _is_docx2pdf_available() -> bool:
    try:
        import docx2pdf  # noqa: F401
        return True
    except Exception:
        return False


def _run_soffice_convert_for_presentation(input_path: str, output_dir: str, timeout: int = 300) -> Tuple[str, str]:
    """
    Run LibreOffice to convert input_path -> PDF in output_dir.
    Uses a per-job UserInstallation to avoid profile locking.
    Returns (stdout, stderr) on success, raises RuntimeError with details on failure.
    """
    user_profile_dir = os.path.join(output_dir, "lo_profile")
    os.makedirs(user_profile_dir, exist_ok=True)

    cmd = [
        "soffice",
        "--headless",
        "--nologo",
        "--nodefault",
        "--norestore",
        "--invisible",
        "--convert-to",
        "pdf:writer_pdf_Export",
        "--outdir",
        output_dir,
        input_path,
        # ✅ SINGLE dash here:
        f"-env:UserInstallation=file://{user_profile_dir}",
    ]

    env = os.environ.copy()
    env["HOME"] = output_dir

    try:
        proc = subprocess.run(
            cmd,
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            timeout=timeout,
            text=True,
            env=env,
        )
        logger.debug("soffice stdout: %s", proc.stdout)
        logger.debug("soffice stderr: %s", proc.stderr)
        return proc.stdout, proc.stderr
    except FileNotFoundError:
        raise RuntimeError("LibreOffice 'soffice' not found on PATH.")
    except subprocess.CalledProcessError as e:
        stdout = getattr(e, "stdout", "") or ""
        stderr = getattr(e, "stderr", "") or ""
        logger.error(
            "soffice conversion failed. exit=%s stdout=%s stderr=%s",
            getattr(e, "returncode", None),
            stdout,
            stderr,
        )
        raise RuntimeError(
            f"LibreOffice conversion failed (exit {getattr(e, 'returncode', 'unknown')}).\n\nstdout:\n{stdout}\n\nstderr:\n{stderr}"
        )
    except subprocess.TimeoutExpired:
        logger.error("soffice conversion timed out after %s seconds for %s", timeout, input_path)
        raise RuntimeError(f"LibreOffice conversion timed out after {timeout} seconds.")

def _run_soffice_convert_with_output(input_path: str, output_dir: str, timeout: int = 300) -> Tuple[str, str]:
    """Alias kept for compatibility (identical behavior)."""
    return _run_soffice_convert_for_presentation(input_path, output_dir, timeout=timeout)


# -------------------------
# PPTX -> PDF
# -------------------------

def pptx_to_pdf_bytes(file_obj):
    """
    Convert PPTX → PDF without LibreOffice.
    Strategy:
      1. Load PPTX via python-pptx.
      2. Render each slide as an image (via pillow screenshot).
      3. Combine slides into a PDF.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from PIL import Image, ImageDraw
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import A4
    from io import BytesIO

    prs = Presentation(file_obj)

    images = []
    # Render slides to simple white-background images
    for slide in prs.slides:
        img = Image.new("RGB", (1600, 900), "white")
        draw = ImageDraw.Draw(img)

        # draw text placeholders only (simple renderer)
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                draw.text((50, 50), shape.text, fill="black")

        images.append(img)

    # Write PDF
    buf = BytesIO()
    pdf = canvas.Canvas(buf, pagesize=A4)
    w, h = A4

    for img in images:
        img_buf = BytesIO()
        img.save(img_buf, "PNG")
        img_buf.seek(0)

        pdf.drawImage(img_buf, 0, 0, width=w, height=h)
        pdf.showPage()

    pdf.save()
    buf.seek(0)
    return buf.read()


# -------------------------
# DOCX -> PDF
# -------------------------

def docx_to_pdf_bytes(file_obj):
    """
    Convert DOCX → PDF without LibreOffice.
    Strategy:
      1. Try docx2pdf (Windows/macOS only).
      2. Fallback: extract text via python-docx and write PDF via reportlab.
    """
    from io import BytesIO
    import tempfile
    import os

    # ---- Try docx2pdf first (if installed & OS supports MS Word)
    try:
        import docx2pdf
        tmp_dir = tempfile.mkdtemp(prefix="docx2pdf_")
        in_path = os.path.join(tmp_dir, "input.docx")
        out_path = os.path.join(tmp_dir, "output.pdf")

        with open(in_path, "wb") as f:
            f.write(file_obj.read())

        try:
            docx2pdf.convert(in_path, out_path)
        except Exception:
            out_path = None

        if out_path and os.path.exists(out_path):
            with open(out_path, "rb") as f:
                return f.read()
    except Exception:
        pass  # docx2pdf not available → fallback

    # ---- Fallback: python-docx + reportlab text-only renderer
    try:
        from docx import Document
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import A4

        file_obj.seek(0)
        doc = Document(file_obj)

        buf = BytesIO()
        pdf = canvas.Canvas(buf, pagesize=A4)
        width, height = A4

        y = height - 40
        for para in doc.paragraphs:
            pdf.drawString(40, y, para.text)
            y -= 16
            if y < 40:
                pdf.showPage()
                y = height - 40

        pdf.save()
        buf.seek(0)
        return buf.read()
    except Exception as e:
        raise RuntimeError(f"DOCX → PDF failed (no LibreOffice): {e}")


# -------------------------
# XLSX -> PDF (soffice or reportlab fallback)
# -------------------------

def _xlsx_to_pdf_reportlab_fallback(file_obj, page_size=A4) -> bytes:
    if not HAS_REPORTLAB:
        raise RuntimeError("ReportLab not available for fallback renderer.")

    try:
        file_obj.seek(0)
    except Exception:
        pass

    try:
        xls = pd.ExcelFile(file_obj)
    except Exception as e:
        raise RuntimeError(f"Failed to read Excel file for fallback renderer: {e}")

    out = BytesIO()
    doc = SimpleDocTemplate(out, pagesize=page_size, rightMargin=36, leftMargin=36, topMargin=36, bottomMargin=36)
    styles = getSampleStyleSheet()
    normal = styles["Normal"]
    normal.fontSize = 8

    flow = []
    for sheet_name in xls.sheet_names:
        try:
            df = xls.parse(sheet_name=sheet_name, header=None, dtype=str)
        except Exception:
            try:
                df = xls.parse(sheet_name=sheet_name)
            except Exception:
                df = None

        if df is None:
            continue

        flow.append(Paragraph(f"<b>Sheet: {sheet_name}</b>", styles["Heading3"]))
        flow.append(Spacer(1, 6))

        table_data = df.fillna("").astype(str).values.tolist()
        max_cols = max((len(r) for r in table_data), default=0)
        if max_cols > 40:
            table_data = [row[:40] for row in table_data]

        if not table_data:
            flow.append(Paragraph("(empty sheet)", normal))
            flow.append(Spacer(1, 12))
        else:
            tbl = Table(table_data, repeatRows=1)
            tbl_style = TableStyle(
                [
                    ("GRID", (0, 0), (-1, -1), 0.25, colors.grey),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ("FONTSIZE", (0, 0), (-1, -1), 8),
                ]
            )
            tbl.setStyle(tbl_style)
            flow.append(tbl)
            flow.append(Spacer(1, 12))

        flow.append(PageBreak())

    if not flow:
        flow.append(Paragraph("Empty Excel file", normal))

    doc.build(flow)
    out.seek(0)
    return out.read()


def xlsx_to_pdf_bytes(file_obj, page_size=A4, backend: str = "soffice", timeout: int = 120) -> bytes:
    try:
        file_obj.seek(0)
    except Exception:
        pass

    # ALWAYS use ReportLab fallback
    return _xlsx_to_pdf_reportlab_fallback(file_obj)

    # If user explicitly asked for reportlab backend or soffice isn't used
    if backend == "reportlab":
        return _xlsx_to_pdf_reportlab_fallback(file_obj, page_size=page_size)

    if HAS_REPORTLAB:
        return _xlsx_to_pdf_reportlab_fallback(file_obj, page_size=page_size)

    raise RuntimeError("No available backend to convert Excel to PDF. Install LibreOffice or ReportLab.")


# -------------------------
# Images -> PDF
# -------------------------

def images_to_pdf_bytes(file_objs: list) -> bytes:
    """
    Convert a list of uploaded image file-like objects to a single multi-page PDF.
    Returns PDF bytes.
    """
    pil_images = []
    for f in file_objs:
        if f is None:
            continue

        if isinstance(f, Image.Image):
            pil_images.append(f.convert("RGB"))
            continue

        try:
            try:
                f.seek(0)
            except Exception:
                pass

            if isinstance(f, (bytes, bytearray)):
                img = Image.open(BytesIO(f))
            else:
                content = f.read()
                if not content:
                    continue
                img = Image.open(BytesIO(content))
        except UnidentifiedImageError:
            continue
        except Exception:
            continue

        try:
            pil_images.append(img.convert("RGB"))
        except Exception:
            try:
                pil_images.append(Image.fromarray(img).convert("RGB"))
            except Exception:
                continue

    if not pil_images:
        raise ValueError("No valid image files were provided.")

    out = BytesIO()
    first, rest = pil_images[0], pil_images[1:]
    first.save(out, format="PDF", save_all=bool(rest), append_images=rest)
    out.seek(0)
    return out.read()


# -------------------------
# PDF -> TXT (with OCR fallback)
# -------------------------

def _extract_with_pdfminer(pdf_bytes: bytes, laparams: Optional[LAParams] = None) -> str:
    if laparams is None:
        laparams = LAParams(char_margin=2.0, line_margin=0.5, word_margin=0.1, boxes_flow=0.5)
    try:
        return pdfminer_extract_text(io.BytesIO(pdf_bytes), laparams=laparams) or ""
    except Exception:
        return ""


def _extract_with_pdfplumber_per_page(pdf_bytes: bytes) -> List[str]:
    pages: List[str] = []
    try:
        with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
            for page in pdf.pages:
                try:
                    txt = page.extract_text(x_tolerance=2, y_tolerance=2) or ""
                except Exception:
                    txt = ""
                pages.append(txt)
    except Exception:
        return []
    return pages


def _ocr_page_image(pil_image: Image.Image, lang: str = "eng") -> str:
    config = "--psm 3 --oem 1"
    try:
        return pytesseract.image_to_string(pil_image, lang=lang, config=config) or ""
    except Exception:
        return ""


def pdf_to_txt_bytes(
    file_obj,
    ocr: bool = True,
    lang: str = "eng",
    join_pages: Optional[str] = None,
    preserve_layout: bool = True,
) -> bytes:
    try:
        file_obj.seek(0)
    except Exception:
        pass

    pdf_bytes = file_obj.read()
    if not pdf_bytes:
        raise ValueError("Empty PDF file.")

    if join_pages is None:
        join_pages = "\n\n----- PAGE BREAK -----\n\n"

    if preserve_layout:
        layout_text = _extract_with_pdfminer(pdf_bytes)
        if layout_text and len(layout_text.strip()) > 50 and "\n" in layout_text:
            full_text = layout_text
            if ocr:
                pages_plumber = _extract_with_pdfplumber_per_page(pdf_bytes)
                if pages_plumber:
                    empty_page_indices = [i for i, t in enumerate(pages_plumber) if len(t.strip()) < 20]
                    if empty_page_indices:
                        pil_pages = convert_from_bytes(pdf_bytes, dpi=300)
                        for idx in empty_page_indices:
                            if idx < len(pil_pages):
                                ocr_text = _ocr_page_image(pil_pages[idx], lang=lang)
                                if ocr_text.strip():
                                    full_text += f"\n\n----- OCR RECOVERY FOR PAGE {idx+1} -----\n\n{ocr_text}"
            return full_text.encode("utf-8")

    pages_text = _extract_with_pdfplumber_per_page(pdf_bytes)
    if (not pages_text or all(len(t.strip()) == 0 for t in pages_text)) and ocr:
        pil_pages = convert_from_bytes(pdf_bytes, dpi=300)
        pages_text = [_ocr_page_image(pil, lang=lang) for pil in pil_pages]

    if ocr and pages_text:
        need_ocr = [i for i, t in enumerate(pages_text) if len(t.strip()) < 20]
        if need_ocr:
            pil_pages = convert_from_bytes(pdf_bytes, dpi=300)
            for idx in need_ocr:
                if idx < len(pil_pages):
                    ocr_text = _ocr_page_image(pil_pages[idx], lang=lang)
                    if len(ocr_text.strip()) > len(pages_text[idx].strip()):
                        pages_text[idx] = ocr_text

    final_text = join_pages.join(pages_text)
    final_text = final_text.replace("\r\n", "\n").replace("\r", "\n")
    return final_text.encode("utf-8")


# -------------------------
# PDF -> PPTX
# -------------------------

def pdf_to_pptx_bytes(file_obj, dpi: int = 150) -> bytes:
    try:
        file_obj.seek(0)
    except Exception:
        pass
    pdf_bytes = file_obj.read()
    if not pdf_bytes:
        raise ValueError("Empty PDF file.")

    images = convert_from_bytes(pdf_bytes, dpi=dpi)
    if not images:
        raise RuntimeError("No pages found in PDF.")

    prs = Presentation()
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    for pil_im in images:
        img_buf = BytesIO()
        pil_im.save(img_buf, format="PNG", optimize=True)
        img_buf.seek(0)

        blank_layout = None
        for layout in prs.slide_layouts:
            if len(layout.placeholders) == 0:
                blank_layout = layout
                break
        if blank_layout is None:
            blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

        slide = prs.slides.add_slide(blank_layout)
        pic = slide.shapes.add_picture(img_buf, 0, 0)
        pic_width = pic.width
        pic_height = pic.height

        scale_w = slide_width / pic_width
        scale_h = slide_height / pic_height
        scale = min(scale_w, scale_h)

        new_width = int(pic_width * scale)
        new_height = int(pic_height * scale)
        left = int((slide_width - new_width) / 2)
        top = int((slide_height - new_height) / 2)

        pic.left = left
        pic.top = top
        pic.width = new_width
        pic.height = new_height

    out = BytesIO()
    prs.save(out)
    out.seek(0)
    return out.read()


# -------------------------
# Helpers
# -------------------------

def _safe_sheet_name(name: str, max_len: int = 31) -> str:
    invalid_chars = ["\\", "/", "*", "[", "]", ":", "?"]
    for ch in invalid_chars:
        name = name.replace(ch, "_")
    name = name.strip()
    if len(name) > max_len:
        name = name[:max_len]
    if not name:
        name = "sheet"
    return name


# -------------------------
# PDF -> EXCEL
# -------------------------

def pdf_to_excel_bytes(file_obj) -> bytes:
    try:
        file_obj.seek(0)
    except Exception:
        pass

    pdf_bytes = file_obj.read()
    if not pdf_bytes:
        raise ValueError("Empty PDF file.")

    with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
        writer_buffer = io.BytesIO()
        with pd.ExcelWriter(writer_buffer, engine="openpyxl") as xlsx_writer:
            any_table = False
            for page_num, page in enumerate(pdf.pages, start=1):
                try:
                    tables = page.extract_tables()
                except Exception:
                    tables = []

                if tables:
                    for t_idx, table in enumerate(tables, start=1):
                        try:
                            df = pd.DataFrame(table)
                        except Exception:
                            df = pd.DataFrame([[" | ".join([str(c) for c in row])] for row in table])

                        sheet_name = _safe_sheet_name(f"page_{page_num}_table_{t_idx}")
                        df.to_excel(xlsx_writer, sheet_name=sheet_name, index=False, header=False)
                        any_table = True

            if not any_table:
                for page_num, page in enumerate(pdf.pages, start=1):
                    text = page.extract_text() or ""
                    rows = text.splitlines() if text else [""]
                    df = pd.DataFrame(rows)
                    sheet_name = _safe_sheet_name(f"page_{page_num}_text")
                    df.to_excel(xlsx_writer, sheet_name=sheet_name, index=False, header=False)

        writer_buffer.seek(0)
        data = writer_buffer.read()

    return data


# -------------------------
# PDF -> DOCX
# -------------------------

def pdf_to_docx_bytes(file_obj):
    """
    Convert PDF → DOCX using pdf2docx only (no LibreOffice).
    """
    from pdf2docx import Converter
    from io import BytesIO
    import tempfile
    import os

    pdf_bytes = file_obj.read()
    if not pdf_bytes:
        raise RuntimeError("Empty PDF")

    tmp = tempfile.mkdtemp(prefix="pdf2docx_")
    in_path = os.path.join(tmp, "input.pdf")
    out_path = os.path.join(tmp, "output.docx")

    with open(in_path, "wb") as f:
        f.write(pdf_bytes)

    conv = Converter(in_path)
    try:
        conv.convert(out_path)
    finally:
        conv.close()

    if not os.path.exists(out_path):
        raise RuntimeError("pdf2docx failed to create DOCX.")

    with open(out_path, "rb") as f:
        return f.read()


# -------------------------
# PDF -> JPG
# -------------------------

def pdf_to_jpg_bytes(file_obj, filename: Optional[str] = None, dpi: int = 200, first_only: bool = False) -> Tuple[bytes, str, str]:
    try:
        file_obj.seek(0)
    except Exception:
        pass
    pdf_bytes = file_obj.read()
    if not filename:
        filename = "output.pdf"
    base_name = filename.rsplit(".", 1)[0]

    images = convert_from_bytes(pdf_bytes, dpi=dpi)
    if not images:
        raise RuntimeError("No pages found in PDF.")

    if first_only or len(images) == 1:
        img_buf = BytesIO()
        images[0].save(img_buf, format="JPEG", quality=85)
        img_bytes = img_buf.getvalue()
        out_name = f"{base_name}.jpg"
        return img_bytes, out_name, "image/jpeg"

    zip_buf = BytesIO()
    with zipfile.ZipFile(zip_buf, mode="w", compression=zipfile.ZIP_DEFLATED) as zf:
        for idx, img in enumerate(images, start=1):
            single_buf = BytesIO()
            img.save(single_buf, format="JPEG", quality=85)
            single_buf.seek(0)
            zf.writestr(f"page_{idx}.jpg", single_buf.read())
    zip_buf.seek(0)
    out_name = f"{base_name}.zip"
    return zip_buf.getvalue(), out_name, "application/zip"